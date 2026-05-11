"""File parsers — turn bytes into a normalized text representation.

Each parser is independent and lazy-imported, so users without (say)
python-docx installed can still ingest plain text and code. Missing parsers
emit a clear warning per file rather than crashing.

The dispatch is by extension, in priority order: tabular formats are
intentionally NOT parsed here — they go through a separate upload path
that triggers the worker's tabular pipeline (Parquet conversion).
"""

from __future__ import annotations

import logging
from dataclasses import dataclass
from pathlib import Path
from typing import Callable, Optional

log = logging.getLogger(__name__)


@dataclass
class ParsedFile:
    """Normalized output from any parser."""

    text: str
    """Canonical UTF-8 prose — the thing Vortex will embed."""

    parser_name: str
    """Which parser produced this — surfaced in metadata for debugging."""

    skipped_reason: Optional[str] = None
    """If non-None, the file was intentionally not parsed."""


# Extensions we route to the tabular pipeline by passing raw bytes to the
# raw bucket (worker auto-detects). The CLI itself doesn't parse them.
TABULAR_EXTENSIONS = {".csv", ".tsv", ".xlsx", ".xls", ".parquet"}

# Extensions to skip outright (binary blobs that nobody wants in RAG).
BINARY_SKIP_EXTENSIONS = {
    ".exe", ".dll", ".so", ".dylib", ".bin", ".o", ".a", ".lib",
    ".class", ".jar", ".pyc", ".pyo", ".whl",
    ".zip", ".tar", ".gz", ".tgz", ".bz2", ".7z", ".rar",
    ".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".webp", ".ico", ".svg",
    ".mp3", ".mp4", ".wav", ".avi", ".mov", ".mkv", ".webm", ".ogg",
    ".woff", ".woff2", ".ttf", ".otf", ".eot",
    ".db", ".sqlite", ".sqlite3",
    ".pdb", ".lock",
    ".DS_Store",
}

# Plain-text extensions — read as UTF-8, no parsing required.
TEXT_EXTENSIONS = {
    ".txt", ".md", ".markdown", ".rst", ".log",
    # Code — embed as-is; semantic search over code is meaningfully useful.
    ".py", ".js", ".ts", ".tsx", ".jsx", ".mjs", ".cjs",
    ".go", ".rs", ".java", ".kt", ".scala", ".clj", ".cljs",
    ".c", ".h", ".cc", ".cpp", ".hpp", ".cxx", ".m", ".mm",
    ".cs", ".vb", ".fs",
    ".rb", ".php", ".pl", ".sh", ".bash", ".zsh", ".fish",
    ".lua", ".dart", ".swift", ".r", ".jl",
    ".sql", ".graphql", ".gql",
    # Config / structured text
    ".json", ".jsonc", ".yaml", ".yml", ".toml", ".ini", ".cfg", ".conf",
    ".env.example", ".gitignore", ".dockerignore",
    ".xml", ".plist",
    ".tf", ".tfvars", ".hcl",
    # Web (raw — for rendered HTML use the HTML parser path below)
    ".css", ".scss", ".sass", ".less",
    # Notebooks (best-effort; full ipynb parsing would extract cell-by-cell)
    ".vue", ".astro",
}


# ─────────────────────────────────────────────────────────────────────────────
# Parsers
# ─────────────────────────────────────────────────────────────────────────────


def parse_text(raw: bytes, path: Path) -> ParsedFile:
    """UTF-8 with replace fallback. Handles code + plain text + config."""
    text = raw.decode("utf-8", errors="replace")
    return ParsedFile(text=text, parser_name="text")


def parse_pdf(raw: bytes, path: Path) -> ParsedFile:
    """PDF via pypdf. For tables, install [docs] to add pdfplumber.

    Pages are joined with double newlines so the chunker can break at them.
    """
    try:
        from io import BytesIO
        from pypdf import PdfReader
    except ImportError:
        return _missing_dep("pypdf", "pip install 'vortex-ingest-cli[docs]'")

    try:
        reader = PdfReader(BytesIO(raw))
    except Exception as e:
        return ParsedFile(text="", parser_name="pdf",
                          skipped_reason=f"unreadable PDF: {e}")

    pages = []
    for i, page in enumerate(reader.pages):
        try:
            pages.append(page.extract_text() or "")
        except Exception as e:
            log.warning("page %d of %s: %s", i, path.name, e)
    text = "\n\n".join(p for p in pages if p.strip())
    return ParsedFile(text=text, parser_name="pdf")


def parse_docx(raw: bytes, path: Path) -> ParsedFile:
    """Word docs via python-docx — preserves paragraph + heading structure."""
    try:
        from io import BytesIO
        from docx import Document
    except ImportError:
        return _missing_dep("python-docx", "pip install 'vortex-ingest-cli[docs]'")

    try:
        doc = Document(BytesIO(raw))
    except Exception as e:
        return ParsedFile(text="", parser_name="docx",
                          skipped_reason=f"unreadable docx: {e}")

    lines = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        # Mark headings so the chunker can use them as split boundaries.
        if para.style and para.style.name.startswith("Heading"):
            lines.append(f"\n## {text}\n")
        else:
            lines.append(text)
    return ParsedFile(text="\n".join(lines), parser_name="docx")


def parse_pptx(raw: bytes, path: Path) -> ParsedFile:
    """PowerPoint slide-by-slide. Each slide → its own section."""
    try:
        from io import BytesIO
        from pptx import Presentation
    except ImportError:
        return _missing_dep("python-pptx", "pip install 'vortex-ingest-cli[docs]'")

    try:
        prs = Presentation(BytesIO(raw))
    except Exception as e:
        return ParsedFile(text="", parser_name="pptx",
                          skipped_reason=f"unreadable pptx: {e}")

    sections = []
    for i, slide in enumerate(prs.slides, 1):
        lines = [f"## Slide {i}"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text.strip())
        sections.append("\n".join(lines))
    return ParsedFile(text="\n\n".join(sections), parser_name="pptx")


def parse_html(raw: bytes, path: Path) -> ParsedFile:
    """HTML via trafilatura — extracts the main content, strips chrome.

    Falls back to raw text if trafilatura isn't installed.
    """
    try:
        import trafilatura
    except ImportError:
        log.warning("trafilatura not installed — falling back to raw HTML "
                    "(install with: pip install 'vortex-ingest-cli[docs]')")
        return parse_text(raw, path)

    try:
        text = trafilatura.extract(raw.decode("utf-8", errors="replace"),
                                   include_links=False, include_tables=True)
    except Exception as e:
        return ParsedFile(text="", parser_name="html",
                          skipped_reason=f"trafilatura failed: {e}")
    if not text:
        return ParsedFile(text="", parser_name="html",
                          skipped_reason="trafilatura found no main content")
    return ParsedFile(text=text, parser_name="html")


def parse_ipynb(raw: bytes, path: Path) -> ParsedFile:
    """Jupyter notebook — concatenate markdown + code cells; drop outputs."""
    import json
    try:
        nb = json.loads(raw.decode("utf-8", errors="replace"))
    except Exception as e:
        return ParsedFile(text="", parser_name="ipynb",
                          skipped_reason=f"unreadable notebook: {e}")

    parts = []
    for cell in nb.get("cells", []):
        ctype = cell.get("cell_type")
        src = cell.get("source", "")
        if isinstance(src, list):
            src = "".join(src)
        if not src.strip():
            continue
        if ctype == "markdown":
            parts.append(src)
        elif ctype == "code":
            parts.append(f"```python\n{src}\n```")
    return ParsedFile(text="\n\n".join(parts), parser_name="ipynb")


def _missing_dep(package: str, install_cmd: str) -> ParsedFile:
    return ParsedFile(
        text="",
        parser_name="missing",
        skipped_reason=f"{package} not installed — run: {install_cmd}",
    )


# ─────────────────────────────────────────────────────────────────────────────
# Dispatch
# ─────────────────────────────────────────────────────────────────────────────

ParserFunc = Callable[[bytes, Path], ParsedFile]

EXT_PARSERS: dict[str, ParserFunc] = {
    ".pdf": parse_pdf,
    ".docx": parse_docx,
    ".pptx": parse_pptx,
    ".html": parse_html,
    ".htm": parse_html,
    ".ipynb": parse_ipynb,
}


def parser_for(path: Path) -> Optional[ParserFunc]:
    """Return the parser callable for `path`, or None if it's text/skip/tabular."""
    ext = path.suffix.lower()
    if ext in EXT_PARSERS:
        return EXT_PARSERS[ext]
    if ext in TEXT_EXTENSIONS or path.name in TEXT_EXTENSIONS:
        return parse_text
    return None


def classify(path: Path) -> str:
    """Return one of: 'text', 'rich', 'tabular', 'binary', 'unknown'.

    Routing:
      - text   → CLI parses, calls vortex_ingest_text
      - rich   → CLI parses with a format-specific parser, calls vortex_ingest_text
      - tabular → CLI uploads raw bytes to raw S3 bucket (worker handles Parquet)
      - binary → skipped
      - unknown → skipped with a warning
    """
    ext = path.suffix.lower()
    if ext in TABULAR_EXTENSIONS:
        return "tabular"
    if ext in BINARY_SKIP_EXTENSIONS or path.name in BINARY_SKIP_EXTENSIONS:
        return "binary"
    if ext in EXT_PARSERS:
        return "rich"
    if ext in TEXT_EXTENSIONS or path.name in TEXT_EXTENSIONS:
        return "text"
    # Files with no extension that look like text (Makefile, Dockerfile, etc.)
    if path.suffix == "" and path.name in {
        "Makefile", "Dockerfile", "Procfile", "Caddyfile",
        "README", "LICENSE", "CHANGELOG", "TODO", "NOTES",
    }:
        return "text"
    return "unknown"
